"""
PowerPoint Generator for Day 1 - Data Modeling Foundations
Creates professional slides with design elements, diagrams, and code blocks
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_shape_with_text(slide, shape_type, left, top, width, height, text, fill_color=None, text_color=None):
    """Add a shape with formatted text"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.text = text
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if text_color:
        shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_code_block(slide, left, top, width, height, code_text):
    """Add a formatted DAX code block"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(43, 43, 43)
    shape.line.color.rgb = RGBColor(200, 200, 200)
    
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.2)
    text_frame.margin_top = Inches(0.1)
    
    p = text_frame.paragraphs[0]
    p.text = code_text
    p.font.name = 'Consolas'
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(230, 230, 230)
    return shape

def create_star_schema_diagram(slide, center_x, center_y):
    """Create visual star schema diagram"""
    # Central fact table
    fact = add_shape_with_text(
        slide, MSO_SHAPE.RECTANGLE, 
        center_x - Inches(0.8), center_y - Inches(0.5), 
        Inches(1.6), Inches(1),
        "SALES\n(Fact)",
        RGBColor(255, 192, 0),
        RGBColor(0, 0, 0)
    )
    fact.text_frame.paragraphs[0].font.bold = True
    
    # Dimension tables
    dimensions = [
        ("CUSTOMERS", center_x - Inches(2.5), center_y - Inches(2)),
        ("PRODUCTS", center_x + Inches(0.9), center_y - Inches(2)),
        ("CALENDAR", center_x - Inches(0.8), center_y + Inches(1.2))
    ]
    
    for dim_name, dim_x, dim_y in dimensions:
        dim = add_shape_with_text(
            slide, MSO_SHAPE.RECTANGLE,
            dim_x, dim_y,
            Inches(1.3), Inches(0.7),
            dim_name,
            RGBColor(68, 114, 196),
            RGBColor(255, 255, 255)
        )
        dim.text_frame.paragraphs[0].font.bold = True
        
        # Add connecting lines
        connector = slide.shapes.add_connector(1, dim_x + Inches(0.65), dim_y + Inches(0.7),
                                               center_x, center_y)
        connector.line.color.rgb = RGBColor(0, 0, 0)
        connector.line.width = Pt(2)

def create_title_slide(prs):
    """Slide 1: Enhanced title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(0, 51, 102)
    banner.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Data Modeling Foundations"
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Building the Backbone of Power BI Solutions"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Day info
    day_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.6))
    day_frame = day_box.text_frame
    day_frame.text = "Day 1 of 12 | Module 3 - Data Modeling (Part 1)"
    day_frame.paragraphs[0].font.size = Pt(18)
    day_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    day_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Power BI branding
    icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.25), Inches(0.3), Inches(1.5), Inches(0.9))
    icon.fill.solid()
    icon.fill.fore_color.rgb = RGBColor(241, 197, 0)
    icon.text = "Power BI"
    icon.text_frame.paragraphs[0].font.size = Pt(20)
    icon.text_frame.paragraphs[0].font.bold = True
    icon.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    icon.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    icon.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def create_content_slide(prs, title_text, content_points):
    """Create content slide with design elements"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title.text_frame.paragraphs[0].font.size = Pt(32)
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(241, 197, 0)
    line.line.fill.background()
    
    content = slide.placeholders[1].text_frame
    content.clear()
    
    for point in content_points:
        p = content.add_paragraph()
        if isinstance(point, tuple):
            p.text = point[0]
            p.level = point[1]
        else:
            p.text = point
            p.level = 0
        p.font.size = Pt(18)
        p.space_before = Pt(6)

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("Creating enhanced PowerPoint with design elements...")
    
    # Slide 1: Title
    create_title_slide(prs)
    print("✓ Slide 1: Enhanced Title")
    
    # Slide 2: Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Today's Agenda"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    content = slide.placeholders[1].text_frame
    content.text = "What We'll Cover"
    for point in ["Data Models", "Fact vs Dimension", "Star Schema", "Relationships", "Common Mistakes", "Practice"]:
        p = content.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
    print("✓ Slide 2: Agenda")
    
    # Content slides
    slides_data = [
        ("Why Data Modeling Matters", [
            "Your data model is 70% of success",
            "", ("Poor modeling = Incorrect calculations", 1),
            ("Poor modeling = Performance issues", 1),
            ("Poor modeling = Maintenance nightmares", 1)
        ]),
        ("Fact Tables", [
            "Contains measurable numbers",
            ("Revenue, Quantity, Cost", 1),
            ("Large and grows over time", 1),
            "", "If you COUNT or SUM it → Fact table"
        ]),
        ("Dimension Tables", [
            "Contains descriptive attributes",
            ("Names, Categories, Dates", 1),
            ("Smaller and stable", 1),
            "", "If you DESCRIBE or FILTER → Dimension"
        ])
    ]
    
    for title, content in slides_data:
        create_content_slide(prs, title, content)
        print(f"✓ Created: {title}")
    
    # Star Schema with diagram
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Star Schema - The Industry Standard"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    create_star_schema_diagram(slide, Inches(5), Inches(3.5))
    
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(3), Inches(3))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = "Central fact table\nSurrounded by dimensions\n\n90% of enterprise\nimplementations"
    for para in text_frame.paragraphs:
        para.font.size = Pt(16)
    print("✓ Created: Star Schema Diagram")
    
    # Dataset with visual tables
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Our Practice Dataset"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    sales_box = add_shape_with_text(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(4.3), Inches(1.2),
        "SALES (Fact)\nOrderID, OrderDate, CustomerID, ProductID\nQuantity, Revenue, Cost",
        RGBColor(255, 192, 0), RGBColor(0, 0, 0)
    )
    sales_box.text_frame.paragraphs[0].font.bold = True
    sales_box.text_frame.paragraphs[0].font.size = Pt(12)
    sales_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    sales_box.text_frame.margin_left = Inches(0.1)
    
    dims = [
        ("CUSTOMERS\nCustomerID, Name, City, Region", Inches(0.5), Inches(3)),
        ("PRODUCTS\nProductID, Name, Category, Price", Inches(5.2), Inches(3)),
        ("CALENDAR\nDate, Year, Quarter, Month, Week", Inches(0.5), Inches(5))
    ]
    
    for text, x, y in dims:
        dim_box = add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.3), Inches(1.2),
                                       text, RGBColor(68, 114, 196), RGBColor(255, 255, 255))
        dim_box.text_frame.paragraphs[0].font.bold = True
        dim_box.text_frame.paragraphs[0].font.size = Pt(12)
        dim_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        dim_box.text_frame.margin_left = Inches(0.1)
    print("✓ Created: Dataset with Visual Tables")
    
    # Filter Flow
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "How Filters Flow Through Relationships"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    prod_box = add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.5), Inches(1.8), Inches(0.8),
                                    "PRODUCTS", RGBColor(68, 114, 196), RGBColor(255, 255, 255))
    prod_box.text_frame.paragraphs[0].font.bold = True
    
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6), Inches(2.5), Inches(0.8), Inches(0.8))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0, 176, 80)
    arrow.text = "Flows"
    arrow.text_frame.paragraphs[0].font.size = Pt(12)
    arrow.text_frame.paragraphs[0].font.bold = True
    arrow.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    sales_box = add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.8), Inches(1.8), Inches(0.8),
                                     "SALES", RGBColor(255, 192, 0), RGBColor(0, 0, 0))
    sales_box.text_frame.paragraphs[0].font.bold = True
    
    rule_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(2))
    rule_frame = rule_box.text_frame
    rule_frame.text = "The Gravity Rule:\n\nFilters flow from\nONE → MANY\n\nDimension → Fact\n(automatic)"
    rule_frame.paragraphs[0].font.size = Pt(18)
    rule_frame.paragraphs[0].font.bold = True
    print("✓ Created: Filter Flow Diagram")
    
    # Summary with DAX
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "With Good Model, DAX is Simple"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    points_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5))
    points_frame = points_box.text_frame
    key_points = [
        "✓ Modeling = 70% of success",
        "✓ Star schema = standard",
        "✓ Filters flow dimension → fact",
        "✓ One-to-many most common",
        "✓ Mark Calendar as date table",
        "✓ Proper model = simple DAX"
    ]
    points_frame.text = key_points[0]
    points_frame.paragraphs[0].font.size = Pt(18)
    for point in key_points[1:]:
        p = points_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.space_after = Pt(10)
    
    code_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.4))
    code_title.text_frame.text = "DAX Examples:"
    code_title.text_frame.paragraphs[0].font.size = Pt(16)
    code_title.text_frame.paragraphs[0].font.bold = True
    
    dax_code = """Total Revenue = 
SUM( Sales[Revenue] )

Revenue YTD = 
TOTALYTD( 
    [Total Revenue], 
    Calendar[Date] 
)

Last Year Sales = 
CALCULATE(
    [Total Revenue],
    SAMEPERIODLASTYEAR(
        Calendar[Date]
    )
)"""
    
    add_code_block(slide, Inches(5.2), Inches(1.8), Inches(4.3), Inches(3.8), dax_code)
    
    note = slide.shapes.add_textbox(Inches(5.2), Inches(5.8), Inches(4.3), Inches(0.8))
    note.text_frame.text = "Works automatically with proper relationships!"
    note.text_frame.paragraphs[0].font.size = Pt(12)
    note.text_frame.paragraphs[0].font.italic = True
    note.text_frame.word_wrap = True
    print("✓ Created: Summary with DAX Code Block")
    
    # Thank you
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Day 1 Complete!"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    action = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(2))
    action.text_frame.text = "Next: Complete exercises & build practice model"
    action.text_frame.paragraphs[0].font.size = Pt(18)
    action.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    print("✓ Created: Thank You")
    
    # Save
    output_file = "PowerBI_Day01_Data_Modeling_Foundations.pptx"
    prs.save(output_file)
    print(f"\n✅ Enhanced PowerPoint created!")
    print(f"📁 File: {output_file}")
    print(f"📊 Slides: {len(prs.slides)}")
    print(f"\n🎨 Features:")
    print(f"  • Professional color scheme")
    print(f"  • Star schema diagram")
    print(f"  • Visual table boxes")
    print(f"  • Filter flow diagram")
    print(f"  • DAX code blocks")

if __name__ == "__main__":
    try:
        main()
    except ImportError:
        print("❌ Error: python-pptx not found")
        print("Install: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error: {e}")
